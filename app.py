#!/usr/bin/env python3
"""
office-tools — 办公效率工具集 Web 应用
"""

import os
import io
import uuid
import time
import shutil
import subprocess
import threading
import concurrent.futures
from pathlib import Path
from PIL import Image as _PIL_Image
from flask import (Flask, render_template, request, send_file,
                   jsonify)

# ── 后台定时清理 ──
_CLEANUP_INTERVAL = 600  # 每10分钟
_last_cleanup = time.time()
_cleanup_lock = threading.Lock()


def _periodic_cleanup():
    """清理超过 30 分钟的旧文件 + 过期 OCR 任务"""
    global _last_cleanup
    now = time.time()
    with _cleanup_lock:
        if now - _last_cleanup < _CLEANUP_INTERVAL:
            return
        _last_cleanup = now
    cutoff = now - 1800  # 30分钟
    for d in [UPLOAD_DIR, OUTPUT_DIR]:
        if not d.exists():
            continue
        for f in d.iterdir():
            try:
                if f.is_file() and f.stat().st_mtime < cutoff:
                    _cleanup(f)
            except Exception:
                pass
    # 清理过期 OCR 任务
    with _ocr_tasks_lock:
        expired = [tid for tid, t in _ocr_tasks.items()
                   if t.get('started_at', 0) < cutoff]
        for tid in expired:
            t = _ocr_tasks.pop(tid, None)
            if t:
                _cleanup(t.get('src_path'))
                _cleanup(t.get('dst_path'))

app = Flask(__name__)
app.config['SECRET_KEY'] = os.urandom(24).hex()
app.config['MAX_CONTENT_LENGTH'] = 500 * 1024 * 1024  # 500MB
# Jinja2 配置：让 {%- ... -%} 块能干净地输出
app.jinja_env.trim_blocks = True
app.jinja_env.lstrip_blocks = True

# ── 站点基础信息 ──
SITE_URL = os.environ.get('SITE_URL', 'https://tools.292029.xyz').rstrip('/')
SITE_NAME = '办公效率工具集'

# ── SEO 元数据配置中心 ──
# 每个页面有独特的 title/description/keywords/path，
# 模板通过 url_for_endpoint 反查此表。
_SEO_META: dict = {
    'index': {
        'title': '办公效率工具集 - 免费在线PDF/图片/编码转换工具',
        'description': '免费的在线办公工具集：PDF 转 Word/合并/拆分/压缩/加水印/加密、图片压缩/格式转换/转 PDF、XLS 转 XLSX/CSV 转换、Base64/JSON/时间戳编码转换。无需注册，保护隐私。',
        'keywords': '在线办公工具,PDF转换,图片处理,免费工具,在线工具,文件转换,编码转换,二维码生成',
        'path': '/',
    },
    'xls_to_xlsx': {
        'title': 'XLS 转 XLSX 在线转换工具 - 免费保留格式与图片',
        'description': '在线将 Excel 97-2003 (.xls) 转换为现代 .xlsx 格式，基于 LibreOffice 引擎，完整保留内容、格式、图片、图表。单文件最大 500MB，30 分钟自动清理。',
        'keywords': 'XLS转XLSX,Excel格式转换,xls,xlsx,Excel 97-2003,LibreOffice,在线转换,文件转换',
        'path': '/xls-to-xlsx',
    },
    'pdf_to_word': {
        'title': 'PDF 转 Word (DOCX) 在线工具 - 保留原始排版与 OCR',
        'description': '在线将 PDF 转换为可编辑 Word (DOCX) 文档，基于 pdf2docx 引擎，完整保留原始排版、表格、字体和图片。支持扫描件 OCR 识别，单文件最大 200MB。',
        'keywords': 'PDF转Word,PDF转DOCX,PDF转换,在线PDF转换,OCR,扫描件识别,文档转换,pdf2docx',
        'path': '/pdf-to-word',
    },
    'image_compress': {
        'title': '图片压缩在线工具 - JPG/PNG/WebP 智能压缩',
        'description': '在线压缩 JPG/PNG/WebP 图片，自定义质量和尺寸，智能保持视觉质量。支持批量上传，单文件最大 50MB，免费无需注册。',
        'keywords': '图片压缩,在线压缩,JPG压缩,PNG压缩,WebP压缩,图片优化,智能压缩,文件压缩',
        'path': '/image-compress',
    },
    'image_convert': {
        'title': '图片格式转换工具 - JPG/PNG/WebP/BMP/GIF/TIFF',
        'description': '在线图片格式转换，支持 JPG/PNG/WebP/BMP/GIF/TIFF 互转。保留原图质量，可设置输出尺寸，单文件最大 50MB。',
        'keywords': '图片格式转换,图片转换,JPG转PNG,PNG转JPG,WebP转换,在线转换,免费图片工具',
        'path': '/image-convert',
    },
    'images_to_pdf': {
        'title': '图片转 PDF 在线工具 - 多图合并 PDF',
        'description': '在线将多张图片（JPG/PNG/WebP）合并为一个 PDF 文件，每张图片一页。可调整顺序，支持 1-10 张图片，单文件最大 50MB。',
        'keywords': '图片转PDF,图片合并,在线PDF制作,JPG转PDF,PNG转PDF,多图合并,PDF生成',
        'path': '/images-to-pdf',
    },
    'hash_check': {
        'title': '文件哈希校验工具 - MD5/SHA1/SHA256/SHA512',
        'description': '在线计算文件 MD5/SHA1/SHA256/SHA512 哈希值，验证文件完整性和安全性。本地计算，文件不上传服务器，保护隐私。',
        'keywords': '哈希校验,文件校验,MD5,SHA1,SHA256,SHA512,文件完整性,本地计算,哈希值',
        'path': '/hash-check',
    },
    'base64': {
        'title': 'Base64 编解码在线工具 - 支持 UTF-8 中文',
        'description': '在线 Base64 编码解码工具，支持 UTF-8 中文字符。文本与 Base64 双向转换，结果可一键复制，完全本地处理无需上传。',
        'keywords': 'Base64编码,Base64解码,Base64在线,UTF-8编码,文本编码,本地处理,免费工具',
        'path': '/base64',
    },
    'json_tool': {
        'title': 'JSON 格式化工具 - 美化/压缩/校验/错误定位',
        'description': '在线 JSON 格式化、压缩、校验、错误定位。支持语法高亮，定位错误行列号，开发者必备工具。完全本地处理，文件不上传。',
        'keywords': 'JSON格式化,JSON美化,JSON压缩,JSON校验,JSON修复,开发者工具,在线JSON',
        'path': '/json-tool',
    },
    'timestamp': {
        'title': 'Unix 时间戳转换工具 - 秒/毫秒双向转换',
        'description': '在线 Unix 时间戳与日期时间互转，支持秒和毫秒自动识别。显示北京时间，可一键复制结果，开发者常用工具。',
        'keywords': '时间戳转换,Unix时间戳,时间戳,日期转换,毫秒转换,北京时区,在线工具',
        'path': '/timestamp',
    },
    'pdf_merge': {
        'title': 'PDF 合并在线工具 - 多文件合并保留书签',
        'description': '在线将多个 PDF 文件合并为一个，支持拖拽调整顺序，保留原始书签和元数据。最多 20 个文件，单文件最大 200MB。',
        'keywords': 'PDF合并,PDF拼接,合并PDF,在线PDF,文件合并,书签保留,多文件合并',
        'path': '/pdf-merge',
    },
    'pdf_split': {
        'title': 'PDF 拆分在线工具 - 按页码范围提取页面',
        'description': '在线从 PDF 中提取指定页面或页码范围，生成新的 PDF 文件。支持单页、多页、范围多种提取模式，保留原始质量。',
        'keywords': 'PDF拆分,PDF分割,PDF提取,按页码提取,PDF页面,在线工具,免费PDF',
        'path': '/pdf-split',
    },
    'pdf_compress': {
        'title': 'PDF 压缩在线工具 - 三档压缩节省空间',
        'description': '在线压缩 PDF 文件，通过图像降采样减小体积，三档压缩级别（轻度/中度/激进）。适合扫描件优化，最高可减少 80% 体积。',
        'keywords': 'PDF压缩,PDF优化,文件压缩,扫描件优化,在线PDF,免费PDF工具,PDF减肥',
        'path': '/pdf-compress',
    },
    'qrcode': {
        'title': '二维码生成工具 - 文本/网址/WiFi/邮箱/电话',
        'description': '在线生成文本、网址、WiFi、邮箱、电话等多种类型二维码。支持 PNG 和 SVG 矢量下载，可自定义颜色和尺寸，免费使用。',
        'keywords': '二维码生成,QR Code,WiFi二维码,网址二维码,SVG二维码,在线工具,免费',
        'path': '/qrcode',
    },
    'pdf_watermark': {
        'title': 'PDF 加水印在线工具 - 平铺文字水印',
        'description': '在线为 PDF 添加平铺文字水印，支持自定义字体、大小、颜色、透明度、旋转角度。保护文档版权，输出文件保留原样可阅读。',
        'keywords': 'PDF水印,加水印,PDF版权,文字水印,平铺水印,PDF保护,在线PDF工具',
        'path': '/pdf-watermark',
    },
    'pdf_encrypt': {
        'title': 'PDF 加密/解密在线工具 - AES 256 强加密',
        'description': '在线为 PDF 设置密码保护或移除已知密码。AES 256 位强加密，可分别设置用户密码（打开）和所有者密码（编辑权限）。',
        'keywords': 'PDF加密,PDF解密,PDF密码,AES加密,文档保护,PDF安全,在线加密',
        'path': '/pdf-encrypt',
    },
    'csv_excel': {
        'title': 'CSV 与 Excel 互转工具 - 智能识别编码解决乱码',
        'description': '在线 CSV 与 Excel (XLSX) 互相转换，自动识别编码（UTF-8/GBK/GB2312），彻底解决中文乱码问题。可自定义分隔符。',
        'keywords': 'CSV转Excel,Excel转CSV,CSV转换,编码识别,中文乱码,UTF-8,GBK,在线转换',
        'path': '/csv-excel',
    },
    'zh_convert': {
        'title': 'Office 文档简繁转换工具 - Word/Excel/PPT 在线互转',
        'description': '在线将 Word (.doc/.docx)、Excel (.xls/.xlsx)、PowerPoint (.ppt/.pptx) 文档进行简体与繁体中文互转，基于 OpenCC 引擎，完整保留原排版、字体、表格、图表。支持 6 种 Office 格式。',
        'keywords': '简繁转换,繁体转换,简体转繁体,繁体转简体,Word简繁,Excel简繁,PPT简繁,OpenCC,在线转换,文档转换',
        'path': '/zh-convert',
    },
}


@app.context_processor
def inject_seo():
    """向所有模板注入 site_url / site_name / seo_meta / 工具列表。"""
    return {
        'site_url': SITE_URL,
        'site_name': SITE_NAME,
        'seo_meta': _SEO_META,
    }

BASE_DIR = Path(__file__).resolve().parent
UPLOAD_DIR = BASE_DIR / 'uploads'
OUTPUT_DIR = BASE_DIR / 'output'

UPLOAD_DIR.mkdir(exist_ok=True)
OUTPUT_DIR.mkdir(exist_ok=True)

ALLOWED_EXT_XLS = {'.xls'}
ALLOWED_EXT_PDF = {'.pdf'}
ALLOWED_EXT_IMAGE = {'.jpg', '.jpeg', '.png', '.webp', '.bmp', '.gif', '.tiff'}
ALLOWED_EXT_OFFICE = {'.doc', '.docx', '.ppt', '.pptx', '.xls', '.xlsx'}
MAX_IMAGES_TO_PDF = 10

# 简繁转换方向：旧格式 → 新格式映射（LibreOffice 预处理）
_ZH_OLD_TO_NEW_EXT = {'.doc': '.docx', '.xls': '.xlsx', '.ppt': '.pptx'}
# 简繁转换方向对应文件后缀
_ZH_DIRECTION_SUFFIX = {'s2t': '_s2t', 't2s': '_t2s'}
# 合法方向
_ZH_VALID_DIRECTIONS = {'s2t', 't2s'}

# ── OCR 异步任务 ──
# 任务存储：task_id -> {status, progress, total, message, error, src_path, dst_path, src_filename, started_at}
_ocr_tasks: dict = {}
_ocr_tasks_lock = threading.Lock()
_OCR_EXECUTOR = concurrent.futures.ThreadPoolExecutor(
    max_workers=1, thread_name_prefix='ocr'
)
# PP-OCRv6 tiny 引擎单例（进程内仅加载一次）
_ocr_engine = None
_ocr_engine_lock = threading.Lock()

# OCR 输入限制
OCR_MAX_FILE_BYTES = 50 * 1024 * 1024  # 50MB
OCR_MAX_PAGES = 100


def _cleanup(path: Path):
    """删除文件或目录"""
    try:
        if path.is_file():
            path.unlink(missing_ok=True)
        elif path.is_dir():
            shutil.rmtree(path, ignore_errors=True)
    except Exception:
        pass


# ── 转换函数 ──────────────────────────────────────────────────

def convert_xls_to_xlsx(src_path: Path, dst_path: Path) -> dict:
    """
    通过 LibreOffice 将 .xls 转换为 .xlsx。
    返回 dict: {"success": bool, "error": str | None}
    """
    os.chmod(src_path, 0o644)
    try:
        work_dir = dst_path.parent.resolve()
        result = subprocess.run(
            [
                'libreoffice',
                '--headless',
                '--norestore',
                '--nofirststartwizard',
                '--convert-to', 'xlsx:Calc MS Excel 2007 XML',
                '--outdir', str(work_dir),
                str(src_path.resolve()),
            ],
            capture_output=True,
            text=True,
            timeout=300,
        )

        src_stem = src_path.stem
        generated = work_dir / f'{src_stem}.xlsx'

        if not generated.exists():
            candidates = list(work_dir.glob('*.xlsx'))
            generated = candidates[0] if candidates else None

        if generated and generated.exists():
            if generated.resolve() != dst_path.resolve():
                shutil.move(str(generated), str(dst_path))
            return {"success": True, "error": None}

        return {
            "success": False,
            "error": f"LibreOffice 未生成目标文件。stderr: {result.stderr}",
        }

    except subprocess.TimeoutExpired:
        return {"success": False, "error": "转换超时（>300s）"}
    except FileNotFoundError:
        return {"success": False, "error": "系统中未找到 libreoffice 命令，请安装 LibreOffice"}
    except Exception as e:
        return {"success": False, "error": str(e)}


def convert_pdf_to_docx(src_path: Path, dst_path: Path) -> dict:
    """
    通过 pdf2docx 将 PDF 转换为 DOCX。
    返回 dict: {"success": bool, "error": str | None}
    """
    try:
        from pdf2docx import Converter

        cv = Converter(str(src_path))
        cv.convert(str(dst_path), start=0, end=None)
        cv.close()

        if dst_path.exists() and dst_path.stat().st_size > 0:
            return {"success": True, "error": None}
        return {"success": False, "error": "转换后文件为空"}

    except ImportError:
        return {"success": False, "error": "缺少 pdf2docx 库，请执行: pip install pdf2docx"}
    except Exception as e:
        return {"success": False, "error": f"转换失败: {e}"}


# ── PDF 合并/拆分/压缩 ─────────────────────────────────

ALLOWED_EXT_PDF_MULTI = {'.pdf'}
MAX_PDF_MERGE_FILES = 20
MAX_PDF_SIZE_BYTES = 200 * 1024 * 1024  # 单 PDF 200MB


def _parse_page_ranges(spec: str, total: int) -> list:
    """
    解析页码范围语法（如 "1-3,5,7-9"）。
    返回按页码升序、去重的页码列表（1-indexed）。
    越界或格式错误时抛出 ValueError。
    """
    if not spec or not spec.strip():
        raise ValueError('页码范围不能为空')

    pages = set()
    for part in spec.split(','):
        part = part.strip()
        if not part:
            continue
        if '-' in part:
            segs = part.split('-', 1)
            try:
                a, b = int(segs[0].strip()), int(segs[1].strip())
            except ValueError:
                raise ValueError(f'无效页码范围: "{part}"')
            if a > b:
                a, b = b, a
            for p in range(a, b + 1):
                if p < 1 or p > total:
                    raise ValueError(f'页码 {p} 超出范围 (1-{total})')
                pages.add(p)
        else:
            try:
                p = int(part)
            except ValueError:
                raise ValueError(f'无效页码: "{part}"')
            if p < 1 or p > total:
                raise ValueError(f'页码 {p} 超出范围 (1-{total})')
            pages.add(p)

    if not pages:
        raise ValueError('未指定任何有效页码')

    return sorted(pages)


def merge_pdfs(src_paths: list, dst_path: Path) -> dict:
    """
    合并多个 PDF 文件为一个。
    src_paths: 源 PDF 路径列表（按顺序）
    """
    try:
        from pypdf import PdfWriter

        writer = PdfWriter()
        for src in src_paths:
            writer.append(str(src))

        with open(dst_path, 'wb') as f:
            writer.write(f)
        writer.close()

        if dst_path.exists() and dst_path.stat().st_size > 0:
            return {"success": True, "error": None}
        return {"success": False, "error": "合并后文件为空"}

    except ImportError:
        return {"success": False, "error": "缺少 pypdf 库，请执行: pip install pypdf"}
    except Exception as e:
        return {"success": False, "error": f"合并失败: {e}"}


def split_pdf(src_path: Path, dst_path: Path, page_spec: str) -> dict:
    """
    从 PDF 中按页码范围提取页面，合并为单个 PDF。
    page_spec: 例如 "1-3,5,7-9"
    """
    try:
        from pypdf import PdfReader, PdfWriter

        reader = PdfReader(str(src_path))
        total = len(reader.pages)

        try:
            selected = _parse_page_ranges(page_spec, total)
        except ValueError as e:
            return {"success": False, "error": str(e)}

        if len(selected) == total:
            # 选中所有页面 → 直接复制
            shutil.copy(str(src_path), str(dst_path))
        else:
            writer = PdfWriter()
            for p in selected:
                writer.add_page(reader.pages[p - 1])
            with open(dst_path, 'wb') as f:
                writer.write(f)
            writer.close()

        if dst_path.exists() and dst_path.stat().st_size > 0:
            return {"success": True, "error": None,
                    "selected_pages": selected, "total_pages": total}
        return {"success": False, "error": "拆分后文件为空"}

    except ImportError:
        return {"success": False, "error": "缺少 pypdf 库，请执行: pip install pypdf"}
    except Exception as e:
        return {"success": False, "error": f"拆分失败: {e}"}


# 压缩档位：(最大边宽, 最大边高, JPEG 质量)
PDF_COMPRESS_LEVELS = {
    'screen':  (1240, 1754, 75),   # 150 DPI 屏幕浏览
    'email':   (827,  1169, 60),   # 100 DPI 邮件附件
    'extreme': (595,  842,  50),   # 72 DPI 极限压缩
}


# ── 二维码生成 ─────────────────────────────────

QR_EC_MAP = {'L': 1, 'M': 0, 'Q': 3, 'H': 2}


# ── PDF 水印 ─────────────────────────────────

def add_pdf_watermark(src_path: Path, dst_path: Path, text: str, 
                      font_size: int = 40, opacity: float = 0.3, 
                      rotation: int = 45) -> dict:
    """
    为 PDF 每页添加平铺文字水印。
    使用 reportlab 生成水印层，pypdf 合并。
    """
    try:
        from reportlab.pdfgen import canvas
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.cidfonts import UnicodeCIDFont
        from pypdf import PdfReader, PdfWriter
        
        # 注册中文字体
        pdfmetrics.registerFont(UnicodeCIDFont('STSong-Light'))
        
        reader = PdfReader(str(src_path))
        writer = PdfWriter()
        
        for page in reader.pages:
            # 获取页面尺寸
            media_box = page.mediabox
            page_width = float(media_box.width)
            page_height = float(media_box.height)
            
            # 创建水印层
            packet = io.BytesIO()
            c = canvas.Canvas(packet, pagesize=(page_width, page_height))
            
            # 设置透明度和字体
            c.saveState()
            c.setFillColorRGB(0.5, 0.5, 0.5, opacity)
            c.setFont('STSong-Light', font_size)
            
            # 计算水印间距
            text_width = c.stringWidth(text, 'STSong-Light', font_size)
            spacing_x = text_width + 150
            spacing_y = font_size + 200
            
            # 平铺水印
            y = -font_size
            while y < page_height + font_size:
                x = -text_width
                while x < page_width + text_width:
                    c.saveState()
                    c.translate(x, y)
                    c.rotate(rotation)
                    c.drawString(0, 0, text)
                    c.restoreState()
                    x += spacing_x
                y += spacing_y
            
            c.restoreState()
            c.save()
            
            # 合并水印
            packet.seek(0)
            watermark_pdf = PdfReader(packet)
            watermark_page = watermark_pdf.pages[0]
            page.merge_page(watermark_page)
            writer.add_page(page)
        
        with open(dst_path, 'wb') as f:
            writer.write(f)
        
        if dst_path.exists() and dst_path.stat().st_size > 0:
            return {"success": True, "error": None}
        return {"success": False, "error": "水印添加后文件为空"}
        
    except ImportError:
        return {"success": False, "error": "缺少依赖库，请执行: pip install reportlab pypdf"}
    except Exception as e:
        return {"success": False, "error": f"添加水印失败: {e}"}


# ── PDF 加密/解密 ─────────────────────────────────

def encrypt_pdf(src_path: Path, dst_path: Path, password: str, 
                owner_password: str = None) -> dict:
    """
    为 PDF 添加密码保护。
    password: 用户密码（打开文件需要）
    owner_password: 所有者密码（可选，用于限制编辑/打印）
    """
    try:
        from pypdf import PdfReader, PdfWriter
        
        reader = PdfReader(str(src_path))
        writer = PdfWriter()
        
        for page in reader.pages:
            writer.add_page(page)
        
        # 添加密码
        if owner_password:
            writer.encrypt(user_password=password, owner_password=owner_password)
        else:
            writer.encrypt(user_password=password)
        
        with open(dst_path, 'wb') as f:
            writer.write(f)
        
        if dst_path.exists() and dst_path.stat().st_size > 0:
            return {"success": True, "error": None}
        return {"success": False, "error": "加密后文件为空"}
        
    except ImportError:
        return {"success": False, "error": "缺少 pypdf 库，请执行: pip install pypdf"}
    except Exception as e:
        return {"success": False, "error": f"加密失败: {e}"}


def decrypt_pdf(src_path: Path, dst_path: Path, password: str) -> dict:
    """
    解密已加密的 PDF。
    """
    try:
        from pypdf import PdfReader, PdfWriter
        
        reader = PdfReader(str(src_path))
        
        # 检查是否已加密
        if not reader.is_encrypted:
            return {"success": False, "error": "该 PDF 未被加密"}
        
        # 尝试解密
        try:
            status = reader.decrypt(password)
            if status == 0:
                return {"success": False, "error": "密码错误，无法解密"}
        except Exception:
            return {"success": False, "error": "密码错误，无法解密"}
        
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)
        
        with open(dst_path, 'wb') as f:
            writer.write(f)
        
        if dst_path.exists() and dst_path.stat().st_size > 0:
            return {"success": True, "error": None}
        return {"success": False, "error": "解密后文件为空"}
        
    except ImportError:
        return {"success": False, "error": "缺少 pypdf 库，请执行: pip install pypdf"}
    except Exception as e:
        return {"success": False, "error": f"解密失败: {e}"}


# ── CSV ↔ Excel ─────────────────────────────────

def convert_csv_to_excel(src_path: Path, dst_path: Path, 
                         encoding: str = 'utf-8', delimiter: str = ',') -> dict:
    """
    将 CSV 转换为 Excel (.xlsx)。
    自动检测编码，支持多种分隔符。
    """
    try:
        import pandas as pd
        
        # 读取 CSV
        df = None
        if encoding == 'auto':
            # 尝试多种编码
            for enc in ['utf-8', 'utf-8-sig', 'gbk', 'gb2312', 'gb18030', 'big5', 'latin-1']:
                try:
                    df = pd.read_csv(str(src_path), encoding=enc, delimiter=delimiter)
                    break
                except (UnicodeDecodeError, UnicodeError):
                    continue
        else:
            try:
                df = pd.read_csv(str(src_path), encoding=encoding, delimiter=delimiter)
            except (UnicodeDecodeError, UnicodeError):
                # 尝试自动检测编码
                for enc in ['utf-8', 'utf-8-sig', 'gbk', 'gb2312', 'gb18030', 'latin-1']:
                    try:
                        df = pd.read_csv(str(src_path), encoding=enc, delimiter=delimiter)
                        break
                    except (UnicodeDecodeError, UnicodeError):
                        continue
        
        if df is None:
            return {"success": False, "error": "无法识别文件编码，请手动指定"}
        
        # 写入 Excel
        df.to_excel(str(dst_path), index=False, engine='openpyxl')
        
        if dst_path.exists() and dst_path.stat().st_size > 0:
            return {"success": True, "error": None, 
                    "rows": len(df), "columns": len(df.columns)}
        return {"success": False, "error": "转换后文件为空"}
        
    except ImportError:
        return {"success": False, "error": "缺少依赖库，请执行: pip install pandas openpyxl"}
    except Exception as e:
        return {"success": False, "error": f"转换失败: {e}"}


def convert_excel_to_csv(src_path: Path, dst_path: Path, 
                         encoding: str = 'utf-8', delimiter: str = ',',
                         sheet_name: str = None) -> dict:
    """
    将 Excel (.xlsx/.xls) 转换为 CSV。
    支持选择工作表，指定编码和分隔符。
    """
    try:
        import pandas as pd
        
        # 读取 Excel
        if sheet_name:
            df = pd.read_excel(str(src_path), sheet_name=sheet_name, engine='openpyxl')
        else:
            df = pd.read_excel(str(src_path), sheet_name=0, engine='openpyxl')
        
        # 写入 CSV
        df.to_csv(str(dst_path), index=False, encoding=encoding, sep=delimiter)
        
        if dst_path.exists() and dst_path.stat().st_size > 0:
            return {"success": True, "error": None,
                    "rows": len(df), "columns": len(df.columns)}
        return {"success": False, "error": "转换后文件为空"}
        
    except ImportError:
        return {"success": False, "error": "缺少依赖库，请执行: pip install pandas openpyxl"}
    except Exception as e:
        return {"success": False, "error": f"转换失败: {e}"}


def generate_qrcode(payload: str, error_level: str = 'M',
                    fg_color: str = '#000000', bg_color: str = '#ffffff',
                    box_size: int = 10, border: int = 2) -> dict:
    """
    生成二维码，返回 PNG (bytes) + SVG (str) + 元数据。
    payload: 要编码的文本
    error_level: 'L'/'M'/'Q'/'H'
    """
    try:
        import qrcode
        from qrcode.constants import (
            ERROR_CORRECT_L, ERROR_CORRECT_M, ERROR_CORRECT_Q, ERROR_CORRECT_H,
        )
        from qrcode.image.svg import SvgPathImage

        ec_map = {
            'L': ERROR_CORRECT_L, 'M': ERROR_CORRECT_M,
            'Q': ERROR_CORRECT_Q, 'H': ERROR_CORRECT_H,
        }
        ec = ec_map.get(error_level.upper(), ERROR_CORRECT_M)

        if not payload:
            return {"success": False, "error": "内容不能为空"}

        if len(payload) > 2000:
            return {"success": False, "error": "内容过长（最大 2000 字符）"}

        qr = qrcode.QRCode(
            version=None,
            error_correction=ec,
            box_size=box_size,
            border=border,
        )
        qr.add_data(payload)
        qr.make(fit=True)
        modules = qr.modules_count

        # 生成 PNG
        img = qr.make_image(fill_color=fg_color, back_color=bg_color)
        png_buf = io.BytesIO()
        img.save(png_buf, 'PNG')
        png_bytes = png_buf.getvalue()

        # 生成 SVG
        svg_factory = SvgPathImage
        img_svg = qr.make_image(image_factory=svg_factory,
                                fill_color=fg_color, back_color=bg_color)
        svg_buf = io.BytesIO()
        img_svg.save(svg_buf)
        svg_str = svg_buf.getvalue().decode('utf-8')

        return {
            "success": True,
            "error": None,
            "png": png_bytes,
            "svg": svg_str,
            "modules": modules,
            "bytes_len": len(payload.encode('utf-8')),
            "char_len": len(payload),
        }

    except ImportError:
        return {"success": False, "error": "缺少 qrcode 库，请执行: pip install qrcode"}
    except Exception as e:
        return {"success": False, "error": f"生成失败: {e}"}


def compress_pdf(src_path: Path, dst_path: Path, level: str = 'screen') -> dict:
    """
    PDF 压缩：图像降采样 + JPEG 重压缩 + stream 压缩。
    level: 'screen' / 'email' / 'extreme'
    """
    try:
        import pikepdf
        from pikepdf import Pdf, PdfImage, Name
        from PIL import Image

        if level not in PDF_COMPRESS_LEVELS:
            return {"success": False,
                    "error": f'无效压缩级别 "{level}"，可选: {", ".join(PDF_COMPRESS_LEVELS.keys())}'}

        max_w, max_h, jpeg_q = PDF_COMPRESS_LEVELS[level]
        n_images_processed = 0

        with Pdf.open(str(src_path)) as pdf:
            for page in pdf.pages:
                for name, obj in list(page.images.items()):
                    try:
                        pim = PdfImage(obj)
                        if pim.image_mask:
                            continue
                        if pim.width <= max_w and pim.height <= max_h:
                            continue
                        pil = pim.as_pil_image()
                        pil.thumbnail((max_w, max_h), Image.LANCZOS)
                        if pil.mode == 'CMYK':
                            pil = pil.convert('RGB')
                        buf = io.BytesIO()
                        pil.save(buf, 'JPEG', quality=jpeg_q)
                        obj.write(buf.getvalue(), filter=Name.DCTDecode)
                        # 更新尺寸元数据
                        obj.Width = pil.size[0]
                        obj.Height = pil.size[1]
                        n_images_processed += 1
                    except Exception:
                        # 单张图像处理失败不影响整体
                        continue

            pdf.save(str(dst_path), compress_streams=True,
                     object_stream_mode=pikepdf.ObjectStreamMode.generate)

        if dst_path.exists() and dst_path.stat().st_size > 0:
            return {"success": True, "error": None,
                    "level": level, "images_processed": n_images_processed}
        return {"success": False, "error": "压缩后文件为空"}

    except ImportError:
        return {"success": False,
                "error": "缺少 pikepdf 库，请执行: pip install pikepdf"}
    except Exception as e:
        return {"success": False, "error": f"压缩失败: {e}"}


# ── OCR (PP-OCRv6 tiny) ──────────────────────────────────────

def _get_ocr_engine():
    """OCR 引擎单例。首次调用时下载并加载模型。

    使用 RapidOCR (ONNX Runtime) 路线：避免 PaddlePaddle 对老 CPU
    （无 AVX2，如 Intel Xeon E5-2690 v2）不兼容的问题。

    模型默认是 PP-OCRv4 (RapidOCR 内置预转换的 ONNX 模型)，
    中英文识别质量与原 PaddleOCR 几乎一致。
    """
    global _ocr_engine
    if _ocr_engine is not None:
        return _ocr_engine
    with _ocr_engine_lock:
        if _ocr_engine is not None:
            return _ocr_engine
        from rapidocr_onnxruntime import RapidOCR
        _ocr_engine = RapidOCR()
    return _ocr_engine


def _update_task(task_id: str, **kwargs):
    """线程安全地更新 OCR 任务状态"""
    with _ocr_tasks_lock:
        t = _ocr_tasks.get(task_id)
        if t:
            t.update(kwargs)


# ── OCR 版式还原算法 ──────────────────────────────────────

def _ocr_group_lines(boxes, texts, y_tol_ratio=0.5):
    """
    把 (box, text) 列表按 Y 中心聚成行。
    同一行的 Y 中心差 < 平均高度 * y_tol_ratio。
    返回 lines = [[(box, text), ...], ...]，每行内按 X 排序。
    box = [x1, y1, x2, y2]
    """
    if not boxes:
        return []
    items = list(zip(boxes, texts))
    # 初始排序：Y 中心、X
    items.sort(key=lambda x: ((x[0][1] + x[0][3]) / 2, x[0][0]))

    lines = []
    current = [items[0]]
    for it in items[1:]:
        y_center = (it[0][1] + it[0][3]) / 2
        avg_h = sum(b[3] - b[1] for b, _ in current) / len(current)
        last_y = (current[-1][0][1] + current[-1][0][3]) / 2
        if abs(y_center - last_y) < max(avg_h * y_tol_ratio, 5):
            current.append(it)
        else:
            current.sort(key=lambda x: x[0][0])
            lines.append(current)
            current = [it]
    if current:
        current.sort(key=lambda x: x[0][0])
        lines.append(current)
    return lines


def _ocr_split_paragraphs(lines, gap_ratio=0.7):
    """
    根据行间 Y 间距把 lines 切分成段落。
    间距 > 平均行高 * gap_ratio → 新段落。
    """
    if not lines:
        return []
    paragraphs = [[lines[0]]]
    for prev, curr in zip(lines, lines[1:]):
        prev_bottom = max(b[3] for b, _ in prev)
        curr_top = min(b[1] for b, _ in curr)
        gap = curr_top - prev_bottom
        prev_h = sum(b[3] - b[1] for b, _ in prev) / len(prev)
        curr_h = sum(b[3] - b[1] for b, _ in curr) / len(curr)
        avg_h = (prev_h + curr_h) / 2
        if avg_h > 0 and gap > avg_h * gap_ratio:
            paragraphs.append([curr])
        else:
            paragraphs[-1].append(curr)
    return paragraphs


def _ocr_format_line(line):
    """
    行内多个文本块按 X 间距智能加空格（保留原排版的横向间距）。
    """
    if len(line) == 1:
        return line[0][1]
    parts = []
    for i, (box, text) in enumerate(line):
        if i > 0:
            prev_x2 = line[i - 1][0][2]
            curr_x1 = box[0]
            gap = curr_x1 - prev_x2
            # 估算字符宽度
            w = box[2] - box[0]
            n = max(len(text), 1)
            char_w = max(w / n, 8)
            # 间距太小 → 1 个空格；间距大 → 多空格
            if gap < char_w * 0.5:
                n_space = 1
            else:
                n_space = max(1, int(gap / char_w))
            parts.append(' ' * n_space)
        parts.append(text)
    return ''.join(parts)


def _ocr_avg_line_height(lines):
    """计算全页平均行高"""
    all_h = [b[3] - b[1] for line in lines for b, _ in line]
    return sum(all_h) / len(all_h) if all_h else 20


import re as _re

_HEADING_PATTERNS = [
    _re.compile(r'第[一二三四五六七八九十百\d]+[页章节]'),          # "第N页/章/节"
    _re.compile(r'【.+】'),                                       # 【XXX】
    _re.compile(r'^([一二三四五六七八九十]+)、'),                 # 一、二、...
    _re.compile(r'^\d+[\.、]\s*\S'),                            # 1. xxx  数字开头
]


def _ocr_classify_paragraph(para_lines, all_lines, para_idx, total_paras):
    """
    综合判断段落类型。
    返回 (style_name, level) 之一：
      - ('Heading 1', None)     一级标题（字号明显大）
      - ('Heading 2', None)     子标题（【XXX】/ "第X页" 模式）
      - ('List Bullet', None)   列表项
    """
    if not para_lines:
        return ('Normal', None)

    first_line = para_lines[0]
    first_text = _ocr_format_line(first_line).strip()
    avg_h = _ocr_avg_line_height(all_lines)
    line_h = sum(b[3] - b[1] for b, _ in first_line) / len(first_line)

    # 1. 字号 > 1.3x 平均 → 可能是标题
    size_big = line_h > avg_h * 1.25

    # 2. 位置判断：是否在页面上方
    all_y_min = min(b[1] for line in all_lines for b, _ in line)
    all_y_max = max(b[3] for line in all_lines for b, _ in line)
    para_y = first_line[0][0][1]
    in_top_quarter = (para_y - all_y_min) < (all_y_max - all_y_min) * 0.3

    # 3. 内容模式判断
    is_h2_pattern = (
        any(p.search(first_text) for p in _HEADING_PATTERNS[:2])
        or bool(_HEADING_PATTERNS[2].match(first_text))  # 一、XX
    )
    is_numbered = bool(_HEADING_PATTERNS[3].match(first_text))

    # 标题优先级
    if size_big and in_top_quarter:
        return ('Heading 1', None)
    if is_h2_pattern:
        return ('Heading 2', None)
    if size_big:
        return ('Heading 1', None)

    # 列表项：段落只有 1 行 + 以 - 1. 开头
    if len(para_lines) == 1 and (
        first_text.startswith(('- ', '•', '·'))
        or is_numbered
        or _HEADING_PATTERNS[2].match(first_text)
    ):
        return ('List Bullet', None)

    # 段落内多行都是列表项格式（- / 数字.）→ 整段是列表
    if len(para_lines) > 1:
        list_count = 0
        for ln in para_lines:
            ln_text = _ocr_format_line(ln).strip()
            if (ln_text.startswith(('- ', '•', '·'))
                    or _HEADING_PATTERNS[3].match(ln_text)
                    or _HEADING_PATTERNS[2].match(ln_text)):
                list_count += 1
        if list_count >= max(2, len(para_lines) * 0.6):
            return ('List Bullet', None)

    return ('Normal', None)


def _ocr_is_table_row(line):
    """
    判断单行是否为表格行：行内 >= 2 个文本块。
    """
    return len(line) >= 2


def _ocr_detect_tables_keep_order(paragraphs):
    """
    在段落列表中查找连续的单行-多列段落，识别为表格。
    返回: (normal_paragraphs, table_groups, tbl_anchor_indices)
      - normal_paragraphs: 普通段落列表
      - table_groups: 表格组列表，每个元素是 [paragraph, ...] 连续表格段落
      - tbl_anchor_indices: 表格的插入位置（在 normal_paras 中的索引，表示该表格插入到此 normal 之后）
    """
    if not paragraphs:
        return [], [], []

    normal = []           # 普通段落
    tables = []           # 表格组
    anchors = []          # 表格插入位置（针对 normal 的索引）
    current_table = []    # 当前累加的表格段落
    current_anchor = -1   # 当前表格的插入位置（最后一个 normal 的索引）

    for para in paragraphs:
        is_table_row = (len(para) == 1 and _ocr_is_table_row(para[0]))
        if is_table_row:
            current_table.append(para)
        else:
            # 结算当前表格
            if current_table:
                col_counts = [len(p[0]) for p in current_table]
                if len(current_table) >= 2 and (max(col_counts) - min(col_counts) <= 1):
                    tables.append(current_table)
                    anchors.append(current_anchor)  # 插入到上一个 normal 后
                else:
                    normal.extend(current_table)
                current_table = []
            normal.append(para)
            current_anchor = len(normal) - 1  # 指向刚刚加入的 normal 索引

    if current_table:
        col_counts = [len(p[0]) for p in current_table]
        if len(current_table) >= 2 and (max(col_counts) - min(col_counts) <= 1):
            tables.append(current_table)
            anchors.append(current_anchor)
        else:
            normal.extend(current_table)

    return normal, tables, anchors


def _ocr_write_table(doc, table_paragraphs):
    """
    把识别出的表格段落写为 docx 表格。
    """
    from docx.shared import Pt
    n_rows = len(table_paragraphs)
    n_cols = max(len(p[0]) for p in table_paragraphs)
    table = doc.add_table(rows=n_rows, cols=n_cols)
    table.style = 'Light Grid Accent 1'
    for r_idx, para in enumerate(table_paragraphs):
        line = para[0]
        for c_idx, (box, text) in enumerate(line):
            if c_idx < n_cols:
                cell = table.rows[r_idx].cells[c_idx]
                cell.text = text
                # 设置单元格中文字体大小
                for paragraph in cell.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(10)


def _ocr_worker(task_id: str, src_path: Path, dst_path: Path):
    """OCR 后台任务（在独立线程中执行）"""
    try:
        _update_task(task_id, status='running', message='正在初始化 OCR 引擎...')
        ocr = _get_ocr_engine()  # 首次会触发模型下载/加载

        from pdf2image import convert_from_path
        from docx import Document

        _update_task(task_id, message='正在解析 PDF 页码...')
        # 先获取总页数（只读取 page count，不渲染）
        from pdf2image.pdf2image import pdfinfo_from_path
        try:
            info = pdfinfo_from_path(str(src_path))
            total_pages = info.get('Pages', 0)
        except Exception:
            total_pages = 0
        if total_pages > OCR_MAX_PAGES:
            _update_task(
                task_id,
                status='failed',
                error=f'PDF 页数 {total_pages} 超过限制 {OCR_MAX_PAGES} 页',
            )
            return
        if total_pages == 0:
            total_pages = 1  # 兜底

        _update_task(task_id, total=total_pages, progress=0,
                     message=f'正在识别 0/{total_pages} 页...')

        images = convert_from_path(str(src_path), dpi=200, fmt='jpeg',
                                   thread_count=2)

        doc = Document()
        # 设置默认中文字体（需本机有对应字体；未找到时回退默认）
        from docx.shared import Pt
        style = doc.styles['Normal']
        style.font.name = 'DejaVu Sans'
        style.font.size = Pt(11)

        for i, img in enumerate(images):
            with _ocr_tasks_lock:
                t = _ocr_tasks.get(task_id)
                if t and t.get('cancel'):
                    _update_task(task_id, status='failed',
                                 error='用户已取消')
                    return

            # 写入临时图片用于 OCR
            tmp_img = dst_path.parent / f'_ocr_{task_id}_p{i}.jpg'
            img.convert('RGB').save(str(tmp_img), 'JPEG', quality=85)
            try:
                result = ocr(str(tmp_img))  # RapidOCR 直接返回 tuple
            finally:
                _cleanup(tmp_img)

            # 提取带坐标的文本块
            # RapidOCR 返回 (detections, scores)
            # detections[i] = [[x1,y1], [x2,y1], [x2,y2], [x1,y2]], text, score
            page_items = []
            if result:
                detections = result[0] if len(result) >= 1 else None
                if detections:
                    for det in detections:
                        if not isinstance(det, (list, tuple)) or len(det) < 2:
                            continue
                        poly = det[0]
                        txt = det[1]
                        if not isinstance(poly, (list, tuple)) or len(poly) < 4:
                            continue
                        if not isinstance(txt, str) or not txt.strip():
                            continue
                        xs = [pt[0] for pt in poly]
                        ys = [pt[1] for pt in poly]
                        box = [min(xs), min(ys), max(xs), max(ys)]
                        page_items.append((box, txt.strip()))

            # 写入 docx
            if i > 0:
                doc.add_page_break()

            if not page_items:
                doc.add_paragraph('')
            else:
                # 1. 按 Y 坐标分行
                boxes = [it[0] for it in page_items]
                texts = [it[1] for it in page_items]
                lines = _ocr_group_lines(boxes, texts)

                # 2. 按行间 Y 间距切段落
                paragraphs = _ocr_split_paragraphs(lines)

                # 2.5 表格检测与切分（保持原顺序）
                normal_paras, table_groups, tbl_anchor = _ocr_detect_tables_keep_order(
                    paragraphs)

                # 3. 按原顺序写入：普通段落 / 表格
                tbl_iter = iter(zip(table_groups, tbl_anchor))
                next_tbl = next(tbl_iter, None)
                for p_idx, para in enumerate(normal_paras):
                    # 检查是否需要在本段前插入表格（anchor == p_idx 表示插到 normal[p_idx] 之前）
                    while next_tbl and next_tbl[1] == p_idx:
                        _ocr_write_table(doc, next_tbl[0])
                        doc.add_paragraph('')
                        next_tbl = next(tbl_iter, None)
                    line_strs = [_ocr_format_line(line) for line in para]
                    text = '\n'.join(line_strs)
                    style, _ = _ocr_classify_paragraph(
                        para, lines, p_idx, len(paragraphs))
                    if style == 'Normal':
                        doc.add_paragraph(text)
                    elif style == 'List Bullet':
                        try:
                            doc.add_paragraph(text, style='List Bullet')
                        except KeyError:
                            doc.add_paragraph(text)
                    else:
                        level = 1 if style == 'Heading 1' else 2
                        try:
                            doc.add_heading(text, level=level)
                        except Exception:
                            doc.add_paragraph(text)
                # 末尾还有表格
                while next_tbl:
                    _ocr_write_table(doc, next_tbl[0])
                    doc.add_paragraph('')
                    next_tbl = next(tbl_iter, None)

            _update_task(
                task_id,
                progress=i + 1,
                message=f'正在识别 {i + 1}/{total_pages} 页...',
            )

        doc.save(str(dst_path))
        _update_task(task_id, status='success',
                     message=f'识别完成，共 {total_pages} 页')

    except ImportError as e:
        _update_task(task_id, status='failed',
                     error=f'缺少依赖库: {e}。请执行: pip install rapidocr_onnxruntime pdf2image Pillow')
    except Exception as e:
        _update_task(task_id, status='failed', error=f'OCR 失败: {e}')


def convert_pdf_ocr_to_docx(src_path: Path, dst_path: Path) -> dict:
    """
    同步入口：实际不直接调用，保留函数签名以兼容 _handle_convert 风格。
    OCR 任务通过 /api/ocr/start 走异步流程。
    """
    return {"success": False, "error": "OCR 任务请通过 /api/ocr/start 提交"}


# ── 图片处理辅助函数 ─────────────────────────────────────

def _image_mime(ext: str) -> str:
    "图片扩展名 → MIME 类型"
    _map = {
        '.jpg': 'image/jpeg', '.jpeg': 'image/jpeg',
        '.png': 'image/png', '.webp': 'image/webp',
        '.bmp': 'image/bmp', '.gif': 'image/gif',
        '.tiff': 'image/tiff',
    }
    return _map.get(ext, 'application/octet-stream')


def _pil_save_image(img, dst_path: Path, ext: str, quality: int = 85):
    """根据扩展名智能保存图片"""
    kwargs = {}
    if ext in ('.jpg', '.jpeg'):
        kwargs.update(quality=quality, optimize=True, subsampling=-1)
    elif ext == '.webp':
        kwargs['quality'] = quality
    elif ext == '.png':
        kwargs['optimize'] = True
    img.save(str(dst_path), **kwargs)


def _pil_prepare_for_jpeg(img, bg_color='white'):
    """将 RGBA/P/LA 模式转 RGB，alpha 用指定颜色填充"""
    if img.mode in ('RGBA', 'P', 'LA'):
        bg = _PIL_Image.new('RGB', img.size, bg_color)
        if img.mode == 'P':
            img = img.convert('RGBA')
        if img.mode == 'RGBA':
            bg.paste(img, mask=img.split()[-1])
        elif img.mode == 'LA':
            bg.paste(img, mask=img.split()[-1])
        return bg
    if img.mode not in ('RGB', 'L'):
        return img.convert('RGB')
    return img



# ── 简繁转换 (Office 文档) ─────────────────────────────────

# OpenCC 转换器单例：s2t (简→繁) / t2s (繁→简)
_zh_converters: dict = {}
_zh_converter_lock = threading.Lock()


def _get_zh_converter(direction: str):
    """获取 OpenCC 转换器单例。direction: 's2t' | 't2s'"""
    if direction in _zh_converters:
        return _zh_converters[direction]
    with _zh_converter_lock:
        if direction not in _zh_converters:
            import opencc
            _zh_converters[direction] = opencc.OpenCC(direction)
    return _zh_converters[direction]


def _zh_convert_docx(src_path: Path, dst_path: Path, convert_fn) -> dict:
    """docx 简繁转换：遍历 paragraphs/tables/headers/footers 的所有 run.text"""
    try:
        from docx import Document
        doc = Document(str(src_path))

        def _conv_paragraphs(paragraphs):
            for p in paragraphs:
                for run in p.runs:
                    if run.text:
                        run.text = convert_fn(run.text)

        _conv_paragraphs(doc.paragraphs)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    _conv_paragraphs(cell.paragraphs)
        for section in doc.sections:
            _conv_paragraphs(section.header.paragraphs)
            _conv_paragraphs(section.footer.paragraphs)
            # 第一页头/尾（如果存在）
            if section.different_first_page_header_footer:
                _conv_paragraphs(section.first_page_header.paragraphs)
                _conv_paragraphs(section.first_page_footer.paragraphs)

        doc.save(str(dst_path))
        return {"success": True, "error": None}
    except ImportError:
        return {"success": False, "error": "缺少 python-docx 库"}
    except Exception as e:
        return {"success": False, "error": f"docx 转换失败: {e}"}


def _zh_convert_xlsx(src_path: Path, dst_path: Path, convert_fn) -> dict:
    """xlsx 简繁转换：遍历所有 sheet 的每个 cell.value（仅字符串）"""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(src_path))
        n_cells = 0
        for ws in wb.worksheets:
            for row in ws.iter_rows():
                for cell in row:
                    if isinstance(cell.value, str) and cell.value:
                        cell.value = convert_fn(cell.value)
                        n_cells += 1
        wb.save(str(dst_path))
        return {"success": True, "error": None, "cells_converted": n_cells}
    except ImportError:
        return {"success": False, "error": "缺少 openpyxl 库"}
    except Exception as e:
        return {"success": False, "error": f"xlsx 转换失败: {e}"}


def _zh_convert_pptx(src_path: Path, dst_path: Path, convert_fn) -> dict:
    """pptx 简繁转换：遍历所有 shapes 的 text_frame + 备注页"""
    try:
        from pptx import Presentation
        prs = Presentation(str(src_path))
        n_runs = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text:
                            run.text = convert_fn(run.text)
                            n_runs += 1
            # 备注页
            if slide.has_notes_slide:
                notes_tf = slide.notes_slide.notes_text_frame
                for para in notes_tf.paragraphs:
                    for run in para.runs:
                        if run.text:
                            run.text = convert_fn(run.text)
                            n_runs += 1
        prs.save(str(dst_path))
        return {"success": True, "error": None, "runs_converted": n_runs}
    except ImportError:
        return {"success": False, "error": "缺少 python-pptx 库"}
    except Exception as e:
        return {"success": False, "error": f"pptx 转换失败: {e}"}


def _zh_preprocess_old_format(src_path: Path, work_dir: Path) -> dict:
    """
    对 .doc/.xls/.ppt 旧格式调用 LibreOffice 转为新格式。
    返回 {"success": bool, "new_path": Path | None, "error": str | None}
    """
    ext = src_path.suffix.lower()
    new_ext = _ZH_OLD_TO_NEW_EXT.get(ext)
    if not new_ext:
        return {"success": False, "new_path": None, "error": f"未知旧格式: {ext}"}

    new_basename = src_path.stem + new_ext
    expected = work_dir / new_basename

    # 如果已存在同名文件先清掉
    if expected.exists():
        _cleanup(expected)

    try:
        result = subprocess.run(
            [
                'libreoffice', '--headless', '--norestore',
                '--nofirststartwizard',
                '--convert-to', new_ext[1:].upper(),  # docx / xlsx / pptx
                '--outdir', str(work_dir),
                str(src_path.resolve()),
            ],
            capture_output=True, text=True, timeout=300,
        )
    except subprocess.TimeoutExpired:
        return {"success": False, "new_path": None, "error": "LibreOffice 预处理超时（>300s）"}
    except FileNotFoundError:
        return {"success": False, "new_path": None, "error": "系统中未找到 libreoffice 命令"}
    except Exception as e:
        return {"success": False, "new_path": None, "error": f"LibreOffice 调用失败: {e}"}

    if expected.exists():
        return {"success": True, "new_path": expected, "error": None}

    # 兜底：扫一下目录里第一个匹配的新格式文件
    candidates = list(work_dir.glob(f'*{new_ext}'))
    if candidates:
        return {"success": True, "new_path": candidates[0], "error": None}

    return {
        "success": False, "new_path": None,
        "error": f"LibreOffice 未生成 {new_ext} 文件。stderr: {result.stderr[:200] if result.stderr else '(empty)'}",
    }


def convert_office_zh(src_path: Path, dst_path: Path, direction: str) -> dict:
    """
    Office 文档简繁转换总入口。
    direction: 's2t' (简→繁) | 't2s' (繁→简)
    旧格式 (.doc/.xls/.ppt) 先用 LibreOffice 转新格式，再做文本转换。
    """
    if direction not in _ZH_VALID_DIRECTIONS:
        return {"success": False, "error": f"无效方向: {direction}"}

    ext = src_path.suffix.lower()
    if ext not in ALLOWED_EXT_OFFICE:
        allowed = ', '.join(sorted(ALLOWED_EXT_OFFICE))
        return {"success": False, "error": f"不支持的文件类型: {ext}，仅支持 {allowed}"}

    try:
        convert_fn = _get_zh_converter(direction).convert
    except Exception as e:
        return {"success": False, "error": f"OpenCC 初始化失败: {e}"}

    # 旧格式：先 LibreOffice 转新格式（产物放在 OUTPUT_DIR）
    work_src = src_path
    converted_from_old = False
    if ext in _ZH_OLD_TO_NEW_EXT:
        os.chmod(src_path, 0o644)
        pre = _zh_preprocess_old_format(src_path, dst_path.parent)
        if not pre['success']:
            return {"success": False, "error": pre['error']}
        work_src = pre['new_path']
        converted_from_old = True

    # 实际新格式扩展名
    new_ext = work_src.suffix.lower()

    # 派发到具体转换器
    if new_ext == '.docx':
        result = _zh_convert_docx(work_src, dst_path, convert_fn)
    elif new_ext == '.xlsx':
        result = _zh_convert_xlsx(work_src, dst_path, convert_fn)
    elif new_ext == '.pptx':
        result = _zh_convert_pptx(work_src, dst_path, convert_fn)
    else:
        return {"success": False, "error": f"内部错误：未实现的扩展名 {new_ext}"}

    # 清理 LibreOffice 中间产物（与 dst_path 同目录）
    if converted_from_old and work_src.resolve() != dst_path.resolve():
        _cleanup(work_src)

    if not result['success']:
        return result

    # 把中间统计也带上
    result['preprocessed'] = converted_from_old
    return result


# ── 通用下载响应 ──────────────────────────────────────────────

def _make_download_response(dst_path, src_filename, new_ext, mime_type, src_path):
    """构造下载响应，下载完成后自动清理临时文件"""
    if not dst_path.exists():
        return jsonify(success=False, error='转换后未找到输出文件'), 500

    resp = send_file(
        str(dst_path),
        as_attachment=True,
        download_name=Path(src_filename).stem + new_ext,
        mimetype=mime_type,
    )

    original_src = src_path
    original_dst = dst_path

    @resp.call_on_close
    def _cleanup_on_close():
        _cleanup(original_src)
        _cleanup(original_dst)

    return resp


# ── 工具通用路由 ──────────────────────────────────────────────

def _handle_convert(ext_set, convert_fn, new_ext, mime_type):
    """通用文件转换处理"""
    if 'file' not in request.files:
        return jsonify(success=False, error='未上传文件'), 400

    file = request.files['file']
    if not file or not file.filename:
        return jsonify(success=False, error='文件名为空'), 400

    ext = Path(file.filename).suffix.lower()
    if ext not in ext_set:
        allowed = ', '.join(ext_set)
        return jsonify(success=False, error=f'不支持的文件类型 "{ext}"，仅支持 {allowed}'), 400

    uid = uuid.uuid4().hex
    src_path = UPLOAD_DIR / f'{uid}{ext}'
    dst_path = OUTPUT_DIR / f'{uid}{new_ext}'

    file.save(str(src_path))

    try:
        result = convert_fn(src_path, dst_path)
        if not result['success']:
            return jsonify(success=False, error=result['error']), 500
        return _make_download_response(dst_path, file.filename, new_ext, mime_type, src_path)
    except Exception as e:
        _cleanup(src_path)
        _cleanup(dst_path)
        return jsonify(success=False, error=str(e)), 500


# ── 页面路由 ──────────────────────────────────────────────────

@app.route('/')
def index():
    return render_template('index.html')


@app.route('/xls-to-xlsx')
def xls_to_xlsx_page():
    return render_template('xls_to_xlsx.html')


@app.route('/pdf-to-word')
def pdf_to_word_page():
    return render_template('pdf_to_word.html')


@app.route('/image-compress')
def image_compress_page():
    return render_template('image_compress.html')


@app.route('/image-convert')
def image_convert_page():
    return render_template('image_convert.html')


@app.route('/images-to-pdf')
def images_to_pdf_page():
    return render_template('images_to_pdf.html')


@app.route('/hash-check')
def hash_check_page():
    return render_template('hash_check.html')


@app.route('/base64')
def base64_page():
    return render_template('base64.html')


@app.route('/json-tool')
def json_tool_page():
    return render_template('json_tool.html')


@app.route('/timestamp')
def timestamp_page():
    return render_template('timestamp.html')


@app.route('/pdf-merge')
def pdf_merge_page():
    return render_template('pdf_merge.html')


@app.route('/pdf-split')
def pdf_split_page():
    return render_template('pdf_split.html')


@app.route('/pdf-compress')
def pdf_compress_page():
    return render_template('pdf_compress.html')


@app.route('/qrcode')
def qrcode_page():
    return render_template('qrcode.html')


@app.route('/pdf-watermark')
def pdf_watermark_page():
    return render_template('pdf_watermark.html')


@app.route('/pdf-encrypt')
def pdf_encrypt_page():
    return render_template('pdf_encrypt.html')


@app.route('/csv-excel')
def csv_excel_page():
    return render_template('csv_excel.html')


@app.route('/zh-convert')
def zh_convert_page():
    return render_template('zh_convert.html')


# ── SEO 基础设施：sitemap.xml + robots.txt ──────────────

@app.route('/sitemap.xml')
def sitemap_xml():
    """动态生成 sitemap.xml, 包含所有页面与 lastmod 时间。"""
    from flask import Response
    lastmod = time.strftime('%Y-%m-%d', time.gmtime())
    urls = []
    for slug, meta in _SEO_META.items():
        urls.append({
            'loc': SITE_URL + meta['path'],
            'lastmod': lastmod,
            'changefreq': 'weekly' if slug == 'index' else 'monthly',
            'priority': '1.0' if slug == 'index' else '0.8',
        })
    body = ['<?xml version="1.0" encoding="UTF-8"?>']
    body.append('<urlset xmlns="http://www.sitemaps.org/schemas/sitemap/0.9">')
    for u in urls:
        body.append('  <url>')
        body.append(f'    <loc>{u["loc"]}</loc>')
        body.append(f'    <lastmod>{u["lastmod"]}</lastmod>')
        body.append(f'    <changefreq>{u["changefreq"]}</changefreq>')
        body.append(f'    <priority>{u["priority"]}</priority>')
        body.append('  </url>')
    body.append('</urlset>')
    return Response('\n'.join(body), mimetype='application/xml; charset=utf-8')


@app.route('/robots.txt')
def robots_txt():
    """动态生成 robots.txt, 引用 sitemap。"""
    from flask import Response
    body = (
        'User-agent: *\n'
        'Allow: /\n'
        'Disallow: /api/\n'
        'Disallow: /uploads/\n'
        'Disallow: /output/\n'
        '\n'
        f'Sitemap: {SITE_URL}/sitemap.xml\n'
    )
    return Response(body, mimetype='text/plain; charset=utf-8')


# ── API 路由 ──────────────────────────────────────────────────

@app.route('/api/convert/xls-to-xlsx', methods=['POST'])
def api_xls_to_xlsx():
    return _handle_convert(ALLOWED_EXT_XLS, convert_xls_to_xlsx, '.xlsx',
                           'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')


@app.route('/api/convert/pdf-to-docx', methods=['POST'])
def api_pdf_to_docx():
    return _handle_convert(ALLOWED_EXT_PDF, convert_pdf_to_docx, '.docx',
                           'application/vnd.openxmlformats-officedocument.wordprocessingml.document')


# ── API：PDF 合并 ────────────────────────────────

@app.route('/api/convert/pdf-merge', methods=['POST'])
def api_pdf_merge():
    files = request.files.getlist('files')
    valid_files = [f for f in files if f and f.filename]

    if not valid_files:
        return jsonify(success=False, error='未上传文件'), 400
    if len(valid_files) < 2:
        return jsonify(success=False, error='至少需要 2 个 PDF 文件'), 400
    if len(valid_files) > MAX_PDF_MERGE_FILES:
        return jsonify(success=False,
                       error=f'最多支持 {MAX_PDF_MERGE_FILES} 个文件'), 400

    uid = uuid.uuid4().hex
    src_paths = []
    try:
        # 保存所有上传文件，先做大小/扩展名校验
        for i, f in enumerate(valid_files):
            ext = Path(f.filename).suffix.lower()
            if ext != '.pdf':
                for sp in src_paths:
                    _cleanup(sp)
                return jsonify(success=False,
                               error=f'文件 "{f.filename}" 不是 PDF 格式'), 400
            # 检查大小
            f.seek(0, os.SEEK_END)
            size = f.tell()
            f.seek(0)
            if size > MAX_PDF_SIZE_BYTES:
                for sp in src_paths:
                    _cleanup(sp)
                return jsonify(success=False,
                               error=f'文件 "{f.filename}" 超过 {MAX_PDF_SIZE_BYTES // 1024 // 1024}MB 限制'), 400

            sp = UPLOAD_DIR / f'{uid}_{i}{ext}'
            f.save(str(sp))
            src_paths.append(sp)

        dst_path = OUTPUT_DIR / f'{uid}.pdf'
        result = merge_pdfs(src_paths, dst_path)
        if not result['success']:
            return jsonify(success=False, error=result['error']), 500

        resp = send_file(
            str(dst_path),
            as_attachment=True,
            download_name=f'merged_{uid[:8]}.pdf',
            mimetype='application/pdf',
        )

        @resp.call_on_close
        def _cleanup_cb():
            for sp in src_paths:
                _cleanup(sp)
            _cleanup(dst_path)

        return resp

    except Exception as e:
        for sp in src_paths:
            _cleanup(sp)
        _cleanup(OUTPUT_DIR / f'{uid}.pdf')
        return jsonify(success=False, error=f'合并失败: {e}'), 500


# ── API：PDF 拆分 ────────────────────────────────

@app.route('/api/convert/pdf-split', methods=['POST'])
def api_pdf_split():
    if 'file' not in request.files:
        return jsonify(success=False, error='未上传文件'), 400
    file = request.files['file']
    if not file or not file.filename:
        return jsonify(success=False, error='文件名为空'), 400
    if Path(file.filename).suffix.lower() != '.pdf':
        return jsonify(success=False, error='仅支持 .pdf 格式'), 400

    page_spec = request.form.get('pages', '').strip()
    if not page_spec:
        return jsonify(success=False, error='请输入要提取的页码范围'), 400

    uid = uuid.uuid4().hex
    src_path = UPLOAD_DIR / f'{uid}.pdf'
    dst_path = OUTPUT_DIR / f'{uid}.pdf'
    file.save(str(src_path))

    try:
        result = split_pdf(src_path, dst_path, page_spec)
        if not result['success']:
            _cleanup(src_path)
            _cleanup(dst_path)
            return jsonify(success=False, error=result['error']), 400

        # 生成友好文件名: 原文件名_pages_1-3,5,7-9.pdf
        # 清理页码串中的特殊字符
        safe_pages = page_spec.replace(',', '_').replace(' ', '')
        original_stem = Path(file.filename).stem
        download_name = f'{original_stem}_pages_{safe_pages}.pdf'

        resp = _make_download_response(
            dst_path, file.filename, '.pdf', 'application/pdf', src_path)

        # 覆盖 download_name（_make_download_response 使用原 stem + new_ext）
        # 这里手动重写 response headers
        resp.headers['Content-Disposition'] = (
            f'attachment; filename="{download_name}"'
        )
        return resp

    except Exception as e:
        _cleanup(src_path)
        _cleanup(dst_path)
        return jsonify(success=False, error=f'拆分失败: {e}'), 500


# ── API：PDF 压缩 ────────────────────────────────

@app.route('/api/convert/pdf-compress', methods=['POST'])
def api_pdf_compress():
    if 'file' not in request.files:
        return jsonify(success=False, error='未上传文件'), 400
    file = request.files['file']
    if not file or not file.filename:
        return jsonify(success=False, error='文件名为空'), 400
    if Path(file.filename).suffix.lower() != '.pdf':
        return jsonify(success=False, error='仅支持 .pdf 格式'), 400

    level = request.form.get('level', 'screen').strip().lower()
    if level not in PDF_COMPRESS_LEVELS:
        return jsonify(success=False,
                       error=f'无效压缩级别 "{level}"，可选: {", ".join(PDF_COMPRESS_LEVELS.keys())}'), 400

    uid = uuid.uuid4().hex
    src_path = UPLOAD_DIR / f'{uid}.pdf'
    dst_path = OUTPUT_DIR / f'{uid}.pdf'
    file.save(str(src_path))

    try:
        orig_size = src_path.stat().st_size
        result = compress_pdf(src_path, dst_path, level)
        if not result['success']:
            _cleanup(src_path)
            _cleanup(dst_path)
            return jsonify(success=False, error=result['error']), 500

        new_size = dst_path.stat().st_size
        # 文件名: 原名_compressed_{level}.pdf
        original_stem = Path(file.filename).stem
        download_name = f'{original_stem}_compressed_{level}.pdf'

        resp = send_file(
            str(dst_path),
            as_attachment=True,
            download_name=download_name,
            mimetype='application/pdf',
        )

        @resp.call_on_close
        def _cleanup_cb():
            _cleanup(src_path)
            _cleanup(dst_path)

        return resp

    except Exception as e:
        _cleanup(src_path)
        _cleanup(dst_path)
        return jsonify(success=False, error=f'压缩失败: {e}'), 500


# ── API：PDF 水印 ────────────────────────────────

@app.route('/api/convert/pdf-watermark', methods=['POST'])
def api_pdf_watermark():
    if 'file' not in request.files:
        return jsonify(success=False, error='未上传文件'), 400
    file = request.files['file']
    if not file or not file.filename:
        return jsonify(success=False, error='文件名为空'), 400
    if Path(file.filename).suffix.lower() != '.pdf':
        return jsonify(success=False, error='仅支持 .pdf 格式'), 400

    text = request.form.get('text', '').strip()
    if not text:
        return jsonify(success=False, error='请输入水印文字'), 400
    if len(text) > 100:
        return jsonify(success=False, error='水印文字过长（最多 100 字符）'), 400

    font_size = max(10, min(200, request.form.get('font_size', 40, type=int)))
    opacity = max(0.05, min(1.0, request.form.get('opacity', 0.3, type=float)))
    rotation = max(-90, min(90, request.form.get('rotation', 45, type=int)))

    uid = uuid.uuid4().hex
    src_path = UPLOAD_DIR / f'{uid}.pdf'
    dst_path = OUTPUT_DIR / f'{uid}.pdf'
    file.save(str(src_path))

    try:
        result = add_pdf_watermark(src_path, dst_path, text, font_size, opacity, rotation)
        if not result['success']:
            _cleanup(src_path)
            _cleanup(dst_path)
            return jsonify(success=False, error=result['error']), 500

        original_stem = Path(file.filename).stem
        download_name = f'{original_stem}_watermarked.pdf'

        resp = send_file(
            str(dst_path),
            as_attachment=True,
            download_name=download_name,
            mimetype='application/pdf',
        )

        @resp.call_on_close
        def _cleanup_cb():
            _cleanup(src_path)
            _cleanup(dst_path)

        return resp

    except Exception as e:
        _cleanup(src_path)
        _cleanup(dst_path)
        return jsonify(success=False, error=f'添加水印失败: {e}'), 500


# ── API：PDF 加密 ────────────────────────────────

@app.route('/api/convert/pdf-encrypt', methods=['POST'])
def api_pdf_encrypt():
    if 'file' not in request.files:
        return jsonify(success=False, error='未上传文件'), 400
    file = request.files['file']
    if not file or not file.filename:
        return jsonify(success=False, error='文件名为空'), 400
    if Path(file.filename).suffix.lower() != '.pdf':
        return jsonify(success=False, error='仅支持 .pdf 格式'), 400

    password = request.form.get('password', '').strip()
    if not password:
        return jsonify(success=False, error='请输入密码'), 400
    if len(password) > 128:
        return jsonify(success=False, error='密码过长（最多 128 字符）'), 400

    owner_password = request.form.get('owner_password', '').strip()
    if owner_password and len(owner_password) > 128:
        return jsonify(success=False, error='所有者密码过长（最多 128 字符）'), 400

    uid = uuid.uuid4().hex
    src_path = UPLOAD_DIR / f'{uid}.pdf'
    dst_path = OUTPUT_DIR / f'{uid}.pdf'
    file.save(str(src_path))

    try:
        result = encrypt_pdf(src_path, dst_path, password, owner_password if owner_password else None)
        if not result['success']:
            _cleanup(src_path)
            _cleanup(dst_path)
            return jsonify(success=False, error=result['error']), 500

        original_stem = Path(file.filename).stem
        download_name = f'{original_stem}_encrypted.pdf'

        resp = send_file(
            str(dst_path),
            as_attachment=True,
            download_name=download_name,
            mimetype='application/pdf',
        )

        @resp.call_on_close
        def _cleanup_cb():
            _cleanup(src_path)
            _cleanup(dst_path)

        return resp

    except Exception as e:
        _cleanup(src_path)
        _cleanup(dst_path)
        return jsonify(success=False, error=f'加密失败: {e}'), 500


# ── API：PDF 解密 ────────────────────────────────

@app.route('/api/convert/pdf-decrypt', methods=['POST'])
def api_pdf_decrypt():
    if 'file' not in request.files:
        return jsonify(success=False, error='未上传文件'), 400
    file = request.files['file']
    if not file or not file.filename:
        return jsonify(success=False, error='文件名为空'), 400
    if Path(file.filename).suffix.lower() != '.pdf':
        return jsonify(success=False, error='仅支持 .pdf 格式'), 400

    password = request.form.get('password', '').strip()
    if not password:
        return jsonify(success=False, error='请输入密码'), 400

    uid = uuid.uuid4().hex
    src_path = UPLOAD_DIR / f'{uid}.pdf'
    dst_path = OUTPUT_DIR / f'{uid}.pdf'
    file.save(str(src_path))

    try:
        result = decrypt_pdf(src_path, dst_path, password)
        if not result['success']:
            _cleanup(src_path)
            _cleanup(dst_path)
            return jsonify(success=False, error=result['error']), 400

        original_stem = Path(file.filename).stem
        download_name = f'{original_stem}_decrypted.pdf'

        resp = send_file(
            str(dst_path),
            as_attachment=True,
            download_name=download_name,
            mimetype='application/pdf',
        )

        @resp.call_on_close
        def _cleanup_cb():
            _cleanup(src_path)
            _cleanup(dst_path)

        return resp

    except Exception as e:
        _cleanup(src_path)
        _cleanup(dst_path)
        return jsonify(success=False, error=f'解密失败: {e}'), 500


# ── API：CSV → Excel ────────────────────────────────

@app.route('/api/convert/csv-to-excel', methods=['POST'])
def api_csv_to_excel():
    if 'file' not in request.files:
        return jsonify(success=False, error='未上传文件'), 400
    file = request.files['file']
    if not file or not file.filename:
        return jsonify(success=False, error='文件名为空'), 400
    if Path(file.filename).suffix.lower() != '.csv':
        return jsonify(success=False, error='仅支持 .csv 格式'), 400

    encoding = request.form.get('encoding', 'utf-8').strip()
    delimiter = request.form.get('delimiter', ',').strip()
    
    # 验证分隔符
    if len(delimiter) != 1:
        return jsonify(success=False, error='分隔符必须是单个字符'), 400

    uid = uuid.uuid4().hex
    src_path = UPLOAD_DIR / f'{uid}.csv'
    dst_path = OUTPUT_DIR / f'{uid}.xlsx'
    file.save(str(src_path))

    try:
        result = convert_csv_to_excel(src_path, dst_path, encoding, delimiter)
        if not result['success']:
            _cleanup(src_path)
            _cleanup(dst_path)
            return jsonify(success=False, error=result['error']), 400

        original_stem = Path(file.filename).stem
        download_name = f'{original_stem}.xlsx'

        resp = send_file(
            str(dst_path),
            as_attachment=True,
            download_name=download_name,
            mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        )

        @resp.call_on_close
        def _cleanup_cb():
            _cleanup(src_path)
            _cleanup(dst_path)

        return resp

    except Exception as e:
        _cleanup(src_path)
        _cleanup(dst_path)
        return jsonify(success=False, error=f'转换失败: {e}'), 500


# ── API：Excel → CSV ────────────────────────────────

@app.route('/api/convert/excel-to-csv', methods=['POST'])
def api_excel_to_csv():
    if 'file' not in request.files:
        return jsonify(success=False, error='未上传文件'), 400
    file = request.files['file']
    if not file or not file.filename:
        return jsonify(success=False, error='文件名为空'), 400
    
    ext = Path(file.filename).suffix.lower()
    if ext not in {'.xlsx', '.xls'}:
        return jsonify(success=False, error='仅支持 .xlsx 或 .xls 格式'), 400

    encoding = request.form.get('encoding', 'utf-8').strip()
    delimiter = request.form.get('delimiter', ',').strip()
    sheet_name = request.form.get('sheet_name', '').strip() or None
    
    # 验证分隔符
    if len(delimiter) != 1:
        return jsonify(success=False, error='分隔符必须是单个字符'), 400

    uid = uuid.uuid4().hex
    src_path = UPLOAD_DIR / f'{uid}{ext}'
    dst_path = OUTPUT_DIR / f'{uid}.csv'
    file.save(str(src_path))

    try:
        result = convert_excel_to_csv(src_path, dst_path, encoding, delimiter, sheet_name)
        if not result['success']:
            _cleanup(src_path)
            _cleanup(dst_path)
            return jsonify(success=False, error=result['error']), 400

        original_stem = Path(file.filename).stem
        download_name = f'{original_stem}.csv'

        resp = send_file(
            str(dst_path),
            as_attachment=True,
            download_name=download_name,
            mimetype='text/csv',
        )

        @resp.call_on_close
        def _cleanup_cb():
            _cleanup(src_path)
            _cleanup(dst_path)

        return resp

    except Exception as e:
        _cleanup(src_path)
        _cleanup(dst_path)
        return jsonify(success=False, error=f'转换失败: {e}'), 500


# ── API：Office 简繁转换 ─────────────────────────────

@app.route('/api/convert/zh-convert', methods=['POST'])
def api_zh_convert():
    """Office 文档简繁转换。direction: s2t (简→繁) | t2s (繁→简)"""
    if 'file' not in request.files:
        return jsonify(success=False, error='未上传文件'), 400
    file = request.files['file']
    if not file or not file.filename:
        return jsonify(success=False, error='文件名为空'), 400

    direction = request.form.get('direction', '').strip().lower()
    if direction not in _ZH_VALID_DIRECTIONS:
        return jsonify(
            success=False,
            error=f'无效方向 "{direction}"，可选: {", ".join(_ZH_VALID_DIRECTIONS)}',
        ), 400

    ext = Path(file.filename).suffix.lower()
    if ext not in ALLOWED_EXT_OFFICE:
        allowed = ', '.join(sorted(ALLOWED_EXT_OFFICE))
        return jsonify(
            success=False,
            error=f'不支持的文件类型 "{ext}"，仅支持 {allowed}',
        ), 400

    uid = uuid.uuid4().hex
    # 实际输出扩展名：旧格式会被转成新格式
    actual_ext = _ZH_OLD_TO_NEW_EXT.get(ext, ext)
    src_path = UPLOAD_DIR / f'{uid}{ext}'
    dst_path = OUTPUT_DIR / f'{uid}{actual_ext}'
    file.save(str(src_path))

    try:
        result = convert_office_zh(src_path, dst_path, direction)
        if not result['success']:
            _cleanup(src_path)
            _cleanup(dst_path)
            return jsonify(success=False, error=result['error']), 500

        if not dst_path.exists() or dst_path.stat().st_size == 0:
            _cleanup(src_path)
            _cleanup(dst_path)
            return jsonify(success=False, error='转换后文件为空'), 500

        # MIME 映射
        mime_map = {
            '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        }
        mime = mime_map.get(actual_ext, 'application/octet-stream')

        # 下载名: 原文件_stem_方向.扩展名
        # 例: 报告.docx + s2t → 报告_s2t.docx
        #     报告.doc + s2t → 报告_s2t.docx （旧格式会自动转新格式）
        original_stem = Path(file.filename).stem
        suffix = _ZH_DIRECTION_SUFFIX[direction]
        download_name = f'{original_stem}{suffix}{actual_ext}'

        resp = send_file(
            str(dst_path),
            as_attachment=True,
            download_name=download_name,
            mimetype=mime,
        )

        @resp.call_on_close
        def _cleanup_cb():
            _cleanup(src_path)
            _cleanup(dst_path)

        return resp

    except Exception as e:
        _cleanup(src_path)
        _cleanup(dst_path)
        return jsonify(success=False, error=f'简繁转换失败: {e}'), 500


# ── API：二维码生成 ─────────────────────────────

@app.route('/api/qrcode', methods=['POST'])
def api_qrcode():
    data = request.get_json(silent=True) or {}
    payload = data.get('payload', '').strip()
    error_level = data.get('error_level', 'M').strip().upper()
    fg_color = data.get('fg_color', '#000000').strip()
    bg_color = data.get('bg_color', '#ffffff').strip()
    box_size = int(data.get('box_size', 10))
    border = int(data.get('border', 2))

    if error_level not in QR_EC_MAP:
        return jsonify(success=False,
                       error=f'无效纠错级别 "{error_level}"'), 400

    if not fg_color.startswith('#') or len(fg_color) not in (4, 7):
        return jsonify(success=False, error='前景色格式错误（需 #RGB 或 #RRGGBB）'), 400
    if not bg_color.startswith('#') or len(bg_color) not in (4, 7):
        return jsonify(success=False, error='背景色格式错误（需 #RGB 或 #RRGGBB）'), 400

    box_size = max(1, min(50, box_size))
    border = max(0, min(10, border))

    result = generate_qrcode(
        payload, error_level, fg_color, bg_color, box_size, border,
    )
    if not result['success']:
        return jsonify(success=False, error=result['error']), 400

    import base64 as _b64
    return jsonify(
        success=True,
        png_base64=_b64.b64encode(result['png']).decode('ascii'),
        svg=result['svg'],
        modules=result['modules'],
        bytes_len=result['bytes_len'],
        char_len=result['char_len'],
    )


# ── API：图片压缩 ────────────────────────────────────────

@app.route('/api/convert/image-compress', methods=['POST'])
def api_image_compress():
    if 'file' not in request.files:
        return jsonify(success=False, error='未上传文件'), 400
    file = request.files['file']
    if not file or not file.filename:
        return jsonify(success=False, error='文件名为空'), 400
    ext = Path(file.filename).suffix.lower()
    if ext not in ALLOWED_EXT_IMAGE:
        allowed = ', '.join(ALLOWED_EXT_IMAGE)
        return jsonify(success=False, error=f'不支持的文件类型 "{ext}"，仅支持 {allowed}'), 400

    quality = max(1, min(100, request.form.get('quality', 80, type=int)))
    max_width = request.form.get('max_width', None, type=int)
    max_height = request.form.get('max_height', None, type=int)

    uid = uuid.uuid4().hex
    src_path = UPLOAD_DIR / f'{uid}{ext}'
    dst_path = OUTPUT_DIR / f'{uid}{ext}'
    file.save(str(src_path))

    try:
        img = _PIL_Image.open(str(src_path))

        # 非 RGB 模式处理（JPEG 不支持 alpha）
        if ext in ('.jpg', '.jpeg'):
            img = _pil_prepare_for_jpeg(img)

        # 缩放
        if max_width or max_height:
            w, h = img.size
            if max_width and w > max_width:
                ratio = max_width / w
                w, h = max_width, int(h * ratio)
            if max_height and h > max_height:
                ratio = max_height / h
                w, h = int(w * ratio), max_height
            img = img.resize((w, h), _PIL_Image.LANCZOS)

        _pil_save_image(img, dst_path, ext, quality)
        img.close()

        if not dst_path.exists() or dst_path.stat().st_size == 0:
            return jsonify(success=False, error='压缩后文件为空'), 500

        return _make_download_response(dst_path, file.filename, ext,
                                       _image_mime(ext), src_path)
    except Exception as e:
        _cleanup(src_path)
        _cleanup(dst_path)
        return jsonify(success=False, error=f'压缩失败: {e}'), 500


# ── API：图片格式互转 ────────────────────────────────────

@app.route('/api/convert/image-format', methods=['POST'])
def api_image_format():
    if 'file' not in request.files:
        return jsonify(success=False, error='未上传文件'), 400
    file = request.files['file']
    if not file or not file.filename:
        return jsonify(success=False, error='文件名为空'), 400
    ext = Path(file.filename).suffix.lower()
    if ext not in ALLOWED_EXT_IMAGE:
        allowed = ', '.join(ALLOWED_EXT_IMAGE)
        return jsonify(success=False, error=f'不支持的文件类型 "{ext}"，仅支持 {allowed}'), 400

    target = request.form.get('target_format', 'jpeg').strip().lower()
    quality = max(1, min(100, request.form.get('quality', 90, type=int)))
    bg_color = request.form.get('bg_color', 'white').strip()

    ext_map = {
        'jpeg': '.jpg', 'jpg': '.jpg',
        'png': '.png', 'webp': '.webp',
        'bmp': '.bmp', 'gif': '.gif',
        'tiff': '.tiff', 'tif': '.tiff',
    }
    if target not in ext_map:
        return jsonify(success=False, error=f'不支持的目标格式 "{target}"'), 400

    dst_ext = ext_map[target]
    uid = uuid.uuid4().hex
    src_path = UPLOAD_DIR / f'{uid}{ext}'
    dst_path = OUTPUT_DIR / f'{uid}{dst_ext}'
    file.save(str(src_path))

    try:
        img = _PIL_Image.open(str(src_path))

        # 透明通道处理
        if dst_ext in ('.jpg', '.bmp', '.gif') and img.mode in ('RGBA', 'P', 'LA'):
            bg = _PIL_Image.new('RGB', img.size, bg_color)
            if img.mode == 'P':
                img = img.convert('RGBA')
            if img.mode == 'RGBA':
                bg.paste(img, mask=img.split()[-1])
            elif img.mode == 'LA':
                bg.paste(img, mask=img.split()[-1])
            img = bg
        elif dst_ext in ('.jpg', '.bmp', '.gif') and img.mode not in ('RGB', 'L'):
            img = img.convert('RGB')

        _pil_save_image(img, dst_path, dst_ext, quality)
        img.close()

        if not dst_path.exists() or dst_path.stat().st_size == 0:
            return jsonify(success=False, error='转换后文件为空'), 500

        return _make_download_response(dst_path, file.filename, dst_ext,
                                       _image_mime(dst_ext), src_path)
    except Exception as e:
        _cleanup(src_path)
        _cleanup(dst_path)
        return jsonify(success=False, error=f'转换失败: {e}'), 500


# ── API：图片转 PDF ──────────────────────────────────────

@app.route('/api/convert/images-to-pdf', methods=['POST'])
def api_images_to_pdf():
    files = request.files.getlist('files')
    if not files or len(files) == 0 or all(f.filename == '' for f in files):
        return jsonify(success=False, error='未上传文件'), 400

    valid_files = [f for f in files if f and f.filename]
    if len(valid_files) > MAX_IMAGES_TO_PDF:
        return jsonify(success=False,
                       error=f'最多支持 {MAX_IMAGES_TO_PDF} 张图片'), 400
    if len(valid_files) < 1:
        return jsonify(success=False, error='未选择有效文件'), 400

    uid = uuid.uuid4().hex
    dst_path = OUTPUT_DIR / f'{uid}.pdf'
    src_paths = []

    try:
        for f in valid_files:
            f_ext = Path(f.filename).suffix.lower()
            if f_ext not in ALLOWED_EXT_IMAGE:
                for sp in src_paths:
                    _cleanup(sp)
                allowed = ', '.join(ALLOWED_EXT_IMAGE)
                return jsonify(success=False,
                               error=f'不支持的文件类型 "{f_ext}"，仅支持 {allowed}'), 400
            sp = UPLOAD_DIR / f'{uid}_{f.filename}'
            f.save(str(sp))
            src_paths.append(sp)

        images = []
        for sp in src_paths:
            img = _PIL_Image.open(str(sp))
            if img.mode != 'RGB':
                img = img.convert('RGB')
            images.append(img)

        if len(images) == 1:
            images[0].save(str(dst_path), 'PDF', resolution=100.0)
        else:
            images[0].save(str(dst_path), 'PDF', save_all=True,
                           append_images=images[1:], resolution=100.0)

        for img in images:
            img.close()

        if not dst_path.exists() or dst_path.stat().st_size == 0:
            return jsonify(success=False, error='生成 PDF 失败'), 500

        resp = send_file(
            str(dst_path),
            as_attachment=True,
            download_name=f'combined_{uid[:8]}.pdf',
            mimetype='application/pdf',
        )

        @resp.call_on_close
        def _cleanup_cb(src_paths_=src_paths, dst_path_=dst_path):
            for sp in src_paths_:
                _cleanup(sp)
            _cleanup(dst_path_)

        return resp

    except Exception as e:
        for sp in src_paths:
            _cleanup(sp)
        _cleanup(dst_path)
        return jsonify(success=False, error=f'转换失败: {e}'), 500

@app.route('/api/ocr/start', methods=['POST'])
def api_ocr_start():
    """提交 OCR 任务，立即返回 task_id"""
    if 'file' not in request.files:
        return jsonify(success=False, error='未上传文件'), 400
    file = request.files['file']
    if not file or not file.filename:
        return jsonify(success=False, error='文件名为空'), 400
    if Path(file.filename).suffix.lower() != '.pdf':
        return jsonify(success=False, error='仅支持 .pdf 格式'), 400

    # 检查文件大小
    file.seek(0, os.SEEK_END)
    size = file.tell()
    file.seek(0)
    if size > OCR_MAX_FILE_BYTES:
        return jsonify(
            success=False,
            error=f'文件大小 {size // 1024 // 1024}MB 超过限制 {OCR_MAX_FILE_BYTES // 1024 // 1024}MB',
        ), 400

    task_id = uuid.uuid4().hex
    src_path = UPLOAD_DIR / f'{task_id}.pdf'
    dst_path = OUTPUT_DIR / f'{task_id}.docx'
    file.save(str(src_path))

    with _ocr_tasks_lock:
        _ocr_tasks[task_id] = {
            'status': 'pending',
            'progress': 0,
            'total': 0,
            'message': '已入队，等待识别...',
            'error': None,
            'src_path': src_path,
            'dst_path': dst_path,
            'src_filename': file.filename,
            'started_at': time.time(),
            'cancel': False,
        }

    _OCR_EXECUTOR.submit(_ocr_worker, task_id, src_path, dst_path)
    return jsonify(success=True, task_id=task_id)


@app.route('/api/ocr/status/<task_id>')
def api_ocr_status(task_id):
    """轮询 OCR 任务状态"""
    with _ocr_tasks_lock:
        t = _ocr_tasks.get(task_id)
    if not t:
        return jsonify(success=False, error='任务不存在或已过期'), 404
    return jsonify(
        success=True,
        status=t['status'],
        progress=t.get('progress', 0),
        total=t.get('total', 0),
        message=t.get('message', ''),
        error=t.get('error'),
    )


@app.route('/api/ocr/download/<task_id>')
def api_ocr_download(task_id):
    """下载 OCR 结果"""
    with _ocr_tasks_lock:
        t = _ocr_tasks.get(task_id)
    if not t:
        return jsonify(success=False, error='任务不存在'), 404
    if t['status'] != 'success':
        return jsonify(success=False, error='任务未完成或已失败'), 400
    dst_path = t['dst_path']
    if not dst_path.exists():
        return jsonify(success=False, error='结果文件已过期'), 404

    resp = send_file(
        str(dst_path),
        as_attachment=True,
        download_name=Path(t['src_filename']).stem + '.docx',
        mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    )

    src_path = t['src_path']
    tid = task_id

    @resp.call_on_close
    def _cleanup_on_close():
        _cleanup(src_path)
        _cleanup(dst_path)
        with _ocr_tasks_lock:
            _ocr_tasks.pop(tid, None)

    return resp


# ── 全局 ──────────────────────────────────────────────────────

@app.after_request
def cleanup_after_request(response):
    """全局：定时清理过期文件"""
    _periodic_cleanup()
    return response


if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=True)
